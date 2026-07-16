#!/usr/bin/env python3
"""Generate a 16:9 Apple/Notion-style portfolio PPTX for DFU Foot Analysis Pipeline."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt


# --- Design tokens ---
BLUE = RGBColor(0x00, 0x71, 0xE3)
BLUE_SOFT = RGBColor(0xE8, 0xF1, 0xFC)
BLUE_DEEP = RGBColor(0x00, 0x58, 0xB0)
TEXT = RGBColor(0x1D, 0x1D, 0x1F)
TEXT_SEC = RGBColor(0x6E, 0x6E, 0x73)
TEXT_MUTED = RGBColor(0x8E, 0x8E, 0x93)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT = RGBColor(0xF5, 0xF7, 0xFA)
BORDER = RGBColor(0xE5, 0xE5, 0xEA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.7)
MARGIN_Y = Inches(0.55)


def set_run(run, *, size=14, bold=False, color=TEXT, name="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_text(shape, text, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return tf


def add_para(tf, text, *, size=14, bold=False, color=TEXT, space_before=0, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or WHITE
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def pill(slide, left, top, width, height, fill, text, text_color=WHITE, size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_text(shape, text, size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def eyebrow(slide, text, top=MARGIN_Y):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, top + Inches(0.08), Inches(0.22), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    box = slide.shapes.add_textbox(MARGIN_X + Inches(0.32), top, Inches(8), Inches(0.3))
    add_text(box, text.upper(), size=11, bold=True, color=BLUE)


def title(slide, text, top=Inches(0.9)):
    box = slide.shapes.add_textbox(MARGIN_X, top, Inches(12), Inches(0.6))
    add_text(box, text, size=30, bold=True, color=TEXT)


def footer(slide, num, total=14):
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.3), Inches(0.3))
    add_text(box, f"{num:02d} / {total:02d}", size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, *, fill=BG_SOFT, line=None):
    return rect(slide, left, top, width, height, fill=fill, line=line)


def metric_card(slide, left, top, width, height, num, cap):
    shape = rect(slide, left, top, width, height, fill=BLUE)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    set_run(run, size=26, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = cap
    set_run(run2, size=11, bold=False, color=WHITE)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # soft top-right wash via a light blue ellipse-like rounded rect
    wash = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(4)
    )
    wash.fill.solid()
    wash.fill.fore_color.rgb = BLUE_SOFT
    wash.line.fill.background()
    # push wash to back by leaving it; content on top is fine for PPT
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 Cover
    s = new_slide(prs)
    brand = s.shapes.add_textbox(MARGIN_X, Inches(1.8), Inches(11), Inches(0.35))
    add_text(brand, "Portfolio · Computer Vision / ML Engineering", size=14, bold=True, color=BLUE)
    h = s.shapes.add_textbox(MARGIN_X, Inches(2.25), Inches(11.5), Inches(1.5))
    tf = add_text(h, "DFU Foot Analysis Pipeline", size=40, bold=True, color=TEXT)
    add_para(tf, "Shared DINOv3 backbone 위에서 발·궤양 세그멘테이션과 DFU 이진 분류를 통합", size=16, color=TEXT_SEC, space_before=12)
    pill(s, MARGIN_X, Inches(4.4), Inches(2.0), Inches(0.38), BLUE, "2026.06 — 2026.07", size=11)
    pill(s, MARGIN_X + Inches(2.15), Inches(4.4), Inches(1.7), Inches(0.38), BG_SOFT, "Solo · 100%", TEXT, 11)
    pill(s, MARGIN_X + Inches(4.0), Inches(4.4), Inches(2.6), Inches(0.38), BG_SOFT, "PyTorch · DINOv3 · Gradio", TEXT, 11)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(5.7), Inches(12), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    meta = s.shapes.add_textbox(MARGIN_X, Inches(5.9), Inches(6), Inches(0.35))
    add_text(meta, "윤수 라 · yoonsu@rexsw.com", size=12, color=TEXT_MUTED)
    link = s.shapes.add_textbox(Inches(7.5), Inches(5.9), Inches(5), Inches(0.35))
    add_text(link, "github.com/yoonsula/dfu-project", size=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    footer(s, 1)

    # 2 Overview
    s = new_slide(prs)
    eyebrow(s, "Project Overview")
    title(s, "프로젝트 설명")
    c1 = card(s, MARGIN_X, Inches(1.7), Inches(5.8), Inches(1.5))
    tf = add_text(c1, "작업 기간", size=11, bold=True, color=BLUE)
    add_para(tf, "약 4주  ·  2026.06.12 — 2026.07.07", size=20, bold=True, color=TEXT, space_before=6)
    add_para(tf, "설계 → 학습 → 추론 UI까지 End-to-End", size=12, color=TEXT_SEC, space_before=6)
    c2 = card(s, MARGIN_X + Inches(6.1), Inches(1.7), Inches(5.8), Inches(1.5))
    tf = add_text(c2, "기여도 / 역할", size=11, bold=True, color=BLUE)
    add_para(tf, "100% · Owner", size=20, bold=True, color=TEXT, space_before=6)
    add_para(tf, "아키텍처 · 데이터 · 학습/평가 · 게이트 · Gradio+WebRTC", size=12, color=TEXT_SEC, space_before=6)
    c3 = card(s, MARGIN_X, Inches(3.4), Inches(11.9), Inches(1.6), fill=WHITE, line=BORDER)
    tf = add_text(c3, "개요", size=11, bold=True, color=BLUE)
    add_para(
        tf,
        "단일 입력 이미지에서 발 영역 탐지 → 궤양 세그멘테이션 → DFU/Other 분류를 수행하는 통합 파이프라인.",
        size=14,
        color=TEXT,
        space_before=8,
    )
    add_para(
        tf,
        "Frozen DINOv3 ViT-S/16을 공유하고 task별 head만 독립 학습·조합하여 학습 효율과 추론 속도를 동시에 확보.",
        size=14,
        color=TEXT_SEC,
        space_before=4,
    )
    for i, (txt, x) in enumerate([
        ("Foot Dice 0.959", 0),
        ("Wound Dice 0.814", 2.3),
        ("DFU Test Acc 96.2%", 4.6),
        ("DFU Test F1 93.1%", 7.1),
    ]):
        pill(s, MARGIN_X + Inches(x), Inches(5.3), Inches(2.15), Inches(0.36), BLUE_SOFT, txt, BLUE_DEEP, 11)
    footer(s, 2)

    # 3 Architecture
    s = new_slide(prs)
    eyebrow(s, "Architecture")
    title(s, "전체 아키텍처")
    nodes = [
        ("1  Input", "512×512 resize\nImageNet normalize", BG_SOFT),
        ("2  DINOv3 ViT-S/16", "Frozen · Feature\n[B,384,H/16,W/16]", BLUE_SOFT),
        ("3  Shared Features", "1회 encode 후\n3개 head 공유", BG_SOFT),
    ]
    for i, (t, sub, fill) in enumerate(nodes):
        left = MARGIN_X + Inches(i * 4.05)
        box = card(s, left, Inches(1.7), Inches(3.7), Inches(1.55), fill=fill, line=BORDER)
        tf = add_text(box, t, size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        for line in sub.split("\n"):
            add_para(tf, line, size=11, color=TEXT_SEC, space_before=4, align=PP_ALIGN.CENTER)
        if i < 2:
            arr = s.shapes.add_textbox(left + Inches(3.55), Inches(2.2), Inches(0.5), Inches(0.4))
            add_text(arr, "→", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    branches = [
        ("SEG · Foot", "FastInstFootHead", "queries=8 · area gate\nratio ∈ [0.08, 0.5]"),
        ("SEG · Wound", "FastInstWoundHead", "queries=16 · foot crop\ncenter ±0.25 guide"),
        ("CLS · DFU", "DFUFeatureHead", "mean pool → MLP\ndfu / other"),
    ]
    for i, (tag, name, desc) in enumerate(branches):
        left = MARGIN_X + Inches(i * 4.05)
        box = card(s, left, Inches(3.55), Inches(3.7), Inches(1.9), fill=WHITE, line=BORDER)
        tf = add_text(box, tag, size=11, bold=True, color=BLUE)
        add_para(tf, name, size=15, bold=True, color=TEXT, space_before=6)
        for line in desc.split("\n"):
            add_para(tf, line, size=12, color=TEXT_SEC, space_before=3)
    note = s.shapes.add_textbox(MARGIN_X, Inches(5.7), Inches(12), Inches(0.4))
    add_text(note, "학습: train.py --task {foot|wound|dfu}  ·  추론: DFUPipelineModel이 checkpoint 조합", size=12, color=TEXT_SEC)
    footer(s, 3)

    # 4 Skills
    s = new_slide(prs)
    eyebrow(s, "Skills")
    title(s, "스킬 정리")
    groups = [
        ("Computer Vision / DL", ["Semantic Segmentation", "Classification", "ViT / DINOv3", "Transfer Learning", "Multi-task Heads"]),
        ("MLOps / Experiment", ["Checkpoint Design", "Early Stopping", "AMP Training", "Hold-out Eval", "Ablation Study"]),
        ("Data Engineering", ["COCO / ImageFolder", "Hard Negatives", "Foot Crop Pipeline", "Class Imbalance", "Multi-source Merge"]),
        ("Productization", ["Inference Gating", "Gradio UI", "WebRTC Realtime", "Capture Guidance", "CLI Tooling"]),
    ]
    for i, (g, chips) in enumerate(groups):
        col = i % 2
        row = i // 2
        left = MARGIN_X + Inches(col * 6.15)
        top = Inches(1.65) + Inches(row * 2.35)
        box = card(s, left, top, Inches(5.9), Inches(2.15), fill=WHITE, line=BORDER)
        tf = add_text(box, g.upper(), size=11, bold=True, color=TEXT_MUTED)
        add_para(tf, "  ·  ".join(chips), size=13, color=TEXT, space_before=10)
    footer(s, 4)

    # 5 Tech stack
    s = new_slide(prs)
    eyebrow(s, "Tech Stack")
    title(s, "사용 기술 스택")
    stacks = [
        ("Core ML", "PyTorch · Torchvision", "학습/추론 루프, AMP, frozen backbone + light head"),
        ("Backbone", "DINOv3 ViT-S/16", "HF Transformers 로컬 스냅샷 · 384-dim feature"),
        ("Heads", "FastInst-style · MLP", "query-based mask · spatial mean pool classifier"),
        ("Serving / UI", "Gradio + FastRTC", "이미지 탭 + WebRTC 실시간 오버레이"),
        ("Data", "COCO · ImageFolder", "Roboflow, FUSeg, AI Hub, Kaggle, PART A"),
        ("Infra", "CUDA · Python 3.10+", "환경변수 경로 · verify_setup · evaluate"),
    ]
    for i, (lab, name, desc) in enumerate(stacks):
        col, row = i % 3, i // 3
        left = MARGIN_X + Inches(col * 4.05)
        top = Inches(1.7) + Inches(row * 2.2)
        box = card(s, left, top, Inches(3.85), Inches(1.95), fill=WHITE if row == 0 else BG_SOFT, line=BORDER if row == 0 else None)
        tf = add_text(box, lab, size=11, bold=True, color=BLUE)
        add_para(tf, name, size=16, bold=True, color=TEXT, space_before=8)
        add_para(tf, desc, size=12, color=TEXT_SEC, space_before=6)
    footer(s, 5)

    # 6 Algorithm shared backbone
    s = new_slide(prs)
    eyebrow(s, "Key Algorithm")
    title(s, "Shared Backbone · 1회 Encode")
    left_box = card(s, MARGIN_X, Inches(1.7), Inches(5.7), Inches(4.4), fill=WHITE, line=BORDER)
    tf = add_text(left_box, "핵심 아이디어", size=12, bold=True, color=BLUE)
    add_para(tf, "추론 시 backbone을 한 번만 실행하고, 동일 feature map을 foot / wound / dfu head가 공유합니다.", size=13, color=TEXT, space_before=10)
    add_para(tf, "왜 중요한가", size=12, bold=True, color=BLUE, space_before=16)
    add_para(tf, "ViT backbone이 연산의 대부분을 차지하므로 multi-task에서도 latency를 최소화하면서 head를 독립적으로 교체·재학습할 수 있습니다.", size=13, color=TEXT_SEC, space_before=6)
    add_para(tf, "backbone 1×  ·  head 독립 checkpoint  ·  task별 재학습", size=12, bold=True, color=BLUE_DEEP, space_before=16)

    code = card(s, MARGIN_X + Inches(5.95), Inches(1.7), Inches(6.0), Inches(4.4), fill=TEXT)
    tf = add_text(code, "# models/pipeline_model.py", size=11, color=TEXT_MUTED)
    for line in [
        "class DFUPipelineModel(nn.Module):",
        "    def encode(self, x):",
        "        return self.backbone(x)  # 1회",
        "",
        "    def predict_foot_logits(self, f, ...):",
        "        return self.foot_head(f)",
        "",
        "    def predict_wound_logits(self, f, ...):",
        "        return self.wound_head(f)",
        "",
        "    def predict_dfu_logits(self, f):",
        "        return self.dfu_head(f)",
    ]:
        add_para(tf, line, size=11, color=WHITE, space_before=2)
    footer(s, 6)

    # 7 Gate + crop
    s = new_slide(prs)
    eyebrow(s, "Key Algorithm")
    title(s, "Gate Logic + Foot Feature Crop")
    g1 = card(s, MARGIN_X, Inches(1.7), Inches(5.7), Inches(2.0), fill=BG_SOFT)
    tf = add_text(g1, "Inference Gate", size=12, bold=True, color=BLUE)
    add_para(tf, "Foot: area ratio ∈ [0.08, 0.5]", size=13, color=TEXT, space_before=8)
    add_para(tf, "Wound: foot 탐지 + 중앙 ±0.25", size=13, color=TEXT, space_before=4)
    add_para(tf, "DFU: foot 탐지 시 실행", size=13, color=TEXT, space_before=4)
    g2 = card(s, MARGIN_X, Inches(3.9), Inches(5.7), Inches(2.0), fill=BG_SOFT)
    tf = add_text(g2, "Wound Crop", size=12, bold=True, color=BLUE)
    add_para(tf, "foot mask bbox를 margin(0.15) 확장한 feature crop만 wound head에 전달 → 배경 노이즈 감소, 국소 정밀도 향상", size=13, color=TEXT, space_before=8)
    code = card(s, MARGIN_X + Inches(5.95), Inches(1.7), Inches(6.0), Inches(4.2), fill=TEXT)
    tf = add_text(code, "# inference/pipeline.py (요약)", size=11, color=TEXT_MUTED)
    for line in [
        "bbox = bbox_from_mask(foot_mask, margin)",
        "crop = features[..., y0:y1, x0:x1]",
        "logits = model.predict_wound_logits(crop)",
        "full[..., y0:y1, x0:x1] = logits",
        "",
        "# DFU head",
        "pooled = features.mean(dim=(2, 3))",
        "return mlp(pooled)  # dfu / other",
    ]:
        add_para(tf, line, size=12, color=WHITE, space_before=3)
    footer(s, 7)

    # 8 Timeline
    s = new_slide(prs)
    eyebrow(s, "Execution Timeline")
    title(s, "프로젝트 진행 타임라인")
    events = [
        ("06.12", "레포 초기화\nbackbone 설계"),
        ("06.16", "세그 학습 루프\nearly stopping"),
        ("06.18", "HF DINOv3\n파이프라인 통합"),
        ("06.22", "Foot crop\n분류 데이터셋"),
        ("06.30", "DFU head 실험\n오분류 리포트"),
        ("07.07", "문서화·평가\nsize 512 확정"),
    ]
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X + Inches(0.4), Inches(2.35), Inches(11.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    for i, (when, what) in enumerate(events):
        left = MARGIN_X + Inches(i * 2.05)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.75), Inches(2.22), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)
        t = s.shapes.add_textbox(left, Inches(2.7), Inches(1.9), Inches(1.2))
        tf = add_text(t, when, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        for line_txt in what.split("\n"):
            add_para(tf, line_txt, size=11, color=TEXT_SEC, space_before=2, align=PP_ALIGN.CENTER)
    metric_card(s, MARGIN_X, Inches(4.5), Inches(3.8), Inches(1.5), "3", "독립 Task Heads")
    metric_card(s, MARGIN_X + Inches(4.05), Inches(4.5), Inches(3.8), Inches(1.5), "1×", "Backbone Encode / Frame")
    metric_card(s, MARGIN_X + Inches(8.1), Inches(4.5), Inches(3.8), Inches(1.5), "E2E", "Train → Infer → UI")
    footer(s, 8)

    def peri_slide(num, eyebrow_txt, title_txt, steps, table_rows=None, extra_cards=None):
        s = new_slide(prs)
        eyebrow(s, eyebrow_txt)
        title(s, title_txt)
        labels = [("P", "문제"), ("E", "실행"), ("R", "성과"), ("I", "인사이트")]
        for i, ((n, lab), body) in enumerate(zip(labels, steps)):
            left = MARGIN_X + Inches(i * 3.05)
            box = card(s, left, Inches(1.65), Inches(2.9), Inches(2.55), fill=BG_SOFT)
            # number badge
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.18), Inches(1.82), Inches(0.32), Inches(0.32))
            badge.adjustments[0] = 0.2
            badge.fill.solid()
            badge.fill.fore_color.rgb = BLUE
            badge.line.fill.background()
            add_text(badge, n, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            tf = box.text_frame
            tf.clear()
            # spacer for badge
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = lab
            set_run(run, size=13, bold=True, color=TEXT)
            add_para(tf, body, size=11, color=TEXT_SEC, space_before=10)
        if table_rows:
            # simple table as text cards / rows
            top = Inches(4.45)
            headers = table_rows[0]
            col_w = 11.9 / len(headers)
            hdr = card(s, MARGIN_X, top, Inches(11.9), Inches(0.4), fill=BG_SOFT, line=BORDER)
            # draw header texts
            for j, h in enumerate(headers):
                tb = s.shapes.add_textbox(MARGIN_X + Inches(j * col_w) + Inches(0.1), top + Inches(0.05), Inches(col_w - 0.1), Inches(0.3))
                add_text(tb, h, size=11, bold=True, color=TEXT)
            for r_i, row in enumerate(table_rows[1:]):
                y = top + Inches(0.42 + r_i * 0.38)
                for j, cell in enumerate(row):
                    tb = s.shapes.add_textbox(MARGIN_X + Inches(j * col_w) + Inches(0.1), y, Inches(col_w - 0.1), Inches(0.35))
                    color = BLUE_DEEP if (r_i == len(table_rows) - 2 and j >= len(row) - 2) else TEXT_SEC
                    add_text(tb, cell, size=11, bold=(color == BLUE_DEEP), color=color)
        if extra_cards:
            for i, (lab, val, desc) in enumerate(extra_cards):
                left = MARGIN_X + Inches(i * 4.05)
                box = card(s, left, Inches(4.5), Inches(3.85), Inches(1.5), fill=BG_SOFT)
                tf = add_text(box, lab, size=11, bold=True, color=BLUE)
                add_para(tf, val, size=18, bold=True, color=TEXT, space_before=6)
                add_para(tf, desc, size=11, color=TEXT_SEC, space_before=4)
        footer(s, num)
        return s

    peri_slide(
        9,
        "Case 01 · Foot Segmentation",
        "발 세그멘테이션",
        [
            "다양한 각도·배경에서 발 영역만 안정 분리, 신체 오탐 감소 필요",
            "Roboflow foot + hard negative · FastInst q=8 · neg oversample",
            "Val Dice 0.959 / IoU 0.935 (foot_head_v2, 512)",
            "Hard negative·empty-mask 가중이 오탐 억제. 분류 crop 품질 기반",
        ],
        table_rows=[
            ["Run", "Image Size", "Best Val Dice", "Best Val IoU"],
            ["foot_head_v1", "512", "0.954", "0.933"],
            ["foot_head_v2", "512", "0.959", "0.935"],
        ],
    )

    peri_slide(
        10,
        "Case 02 · Wound Segmentation",
        "궤양 세그멘테이션",
        [
            "궤양이 작고 경계가 모호, full-image 학습 시 배경에 민감",
            "FUSeg + Wound Image Dataset · q=16 · foot bbox feature crop",
            "Val Dice 0.814 / IoU 0.738 (wound_head_v1)",
            "추론 ROI crop + gate로 실사용 정밀도·UX 동시 개선",
        ],
        extra_cards=[
            ("데이터", "FUSeg + WID", "positive / negative 혼합"),
            ("Crop Margin", "0.15", "infer.py 기본값"),
            ("Guide Gate", "±0.25", "화면 중앙 정렬 조건"),
        ],
    )

    peri_slide(
        11,
        "Case 03 · DFU Classification",
        "당뇨발 이진 분류",
        [
            "출처·구도가 다른 데이터에서 일반화, 클래스 불균형",
            "AI Hub+Kaggle+PART A · foot crop · Linear/MLP ablation",
            "Test Acc 96.2% / F1 93.1% (mlp + aug)",
            "Foot crop이 도메인 갭 축소. MLP+aug가 F1 개선",
        ],
        table_rows=[
            ["Run", "Head", "Val Acc", "Val F1", "Test Acc", "Test F1"],
            ["dfu_head_v1", "linear", "95.1%", "90.9%", "—", "—"],
            ["dfu_head_mlp_1e3", "mlp", "98.0%", "93.4%", "96.1%", "92.9%"],
            ["dfu_head_mlp_1e3_aug", "mlp+aug", "97.8%", "94.4%", "96.2%", "93.1%"],
        ],
    )

    # 12 Performance
    s = new_slide(prs)
    eyebrow(s, "Results")
    title(s, "기능 · 성능 테스트 요약")
    metrics = [("0.959", "Foot Val Dice"), ("0.814", "Wound Val Dice"), ("96.2%", "DFU Test Acc"), ("93.1%", "DFU Test F1")]
    for i, (n, c) in enumerate(metrics):
        metric_card(s, MARGIN_X + Inches(i * 3.05), Inches(1.65), Inches(2.9), Inches(1.25), n, c)
    rows = [
        ("통합 Inference CLI", "infer.py · mask/overlay/JSON", "foot→wound→dfu 게이트 동작"),
        ("Hold-out Evaluation", "scripts/evaluate.py", "dice/iou · acc/P/R/f1"),
        ("DFU Head 비교", "evaluate_dfu_heads_on_test.py", "MLP+aug Test F1 best"),
        ("Realtime UI", "Gradio + FastRTC", "오버레이 + 0.5s JSON"),
        ("촬영 가이드", "area / center gate", "잘못된 입력 사전 차단"),
    ]
    hdr = card(s, MARGIN_X, Inches(3.15), Inches(11.9), Inches(0.4), fill=BG_SOFT, line=BORDER)
    for j, h in enumerate(["기능", "검증 방법", "결과 / 지표"]):
        tb = s.shapes.add_textbox(MARGIN_X + Inches(0.15 + j * 4.0), Inches(3.2), Inches(3.8), Inches(0.3))
        add_text(tb, h, size=11, bold=True, color=TEXT)
    for i, (a, b, c) in enumerate(rows):
        y = Inches(3.6 + i * 0.45)
        for j, cell in enumerate([a, b, c]):
            tb = s.shapes.add_textbox(MARGIN_X + Inches(0.15 + j * 4.0), y, Inches(3.8), Inches(0.4))
            add_text(tb, cell, size=12, color=TEXT_SEC)
    footer(s, 12)

    # 13 GitHub
    s = new_slide(prs)
    eyebrow(s, "Links & Deliverables")
    title(s, "GitHub · 산출물")
    box = card(s, MARGIN_X, Inches(1.7), Inches(11.9), Inches(1.6), fill=WHITE, line=BORDER)
    tf = add_text(box, "Source Repository", size=11, bold=True, color=BLUE)
    add_para(tf, "yoonsula / dfu-project", size=22, bold=True, color=TEXT, space_before=8)
    add_para(tf, "https://github.com/yoonsula/dfu-project", size=14, bold=True, color=BLUE, space_before=6)
    for i, (lab, name, desc) in enumerate([
        ("Training", "train.py", "task dispatcher · foot / wound / dfu"),
        ("Inference", "infer.py", "gated pipeline · timing JSON · overlay"),
        ("Demo UI", "app_gradio.py", "이미지 분석 + WebRTC 실시간"),
    ]):
        left = MARGIN_X + Inches(i * 4.05)
        c = card(s, left, Inches(3.7), Inches(3.85), Inches(1.9), fill=BG_SOFT)
        tf = add_text(c, lab, size=11, bold=True, color=BLUE)
        add_para(tf, name, size=16, bold=True, color=TEXT, space_before=8)
        add_para(tf, desc, size=12, color=TEXT_SEC, space_before=6)
    footer(s, 13)

    # 14 One-liner
    s = new_slide(prs)
    eyebrow(s, "One-liner")
    q = s.shapes.add_textbox(MARGIN_X, Inches(2.0), Inches(12), Inches(3))
    tf = add_text(q, "하나의 frozen backbone으로", size=26, bold=True, color=TEXT)
    add_para(tf, "발 탐지 · 궤양 분할 · DFU 분류를 묶고,", size=26, bold=True, color=TEXT, space_before=8)
    add_para(tf, "게이트와 crop으로 실사용 정밀도까지 끌어올린", size=26, bold=True, color=TEXT, space_before=8)
    add_para(tf, "의료 영상 분석 파이프라인.", size=26, bold=True, color=BLUE, space_before=8)
    pill(s, MARGIN_X, Inches(5.3), Inches(2.4), Inches(0.4), BLUE, "Dice 0.959 / 0.814", WHITE, 12)
    pill(s, MARGIN_X + Inches(2.55), Inches(5.3), Inches(2.0), Inches(0.4), BLUE, "DFU F1 93.1%", WHITE, 12)
    pill(s, MARGIN_X + Inches(4.7), Inches(5.3), Inches(3.8), Inches(0.4), BG_SOFT, "github.com/yoonsula/dfu-project", TEXT, 11)
    foot = s.shapes.add_textbox(MARGIN_X, Inches(6.0), Inches(10), Inches(0.35))
    add_text(foot, "윤수 라 · DFU Foot Analysis Pipeline Portfolio · 2026", size=12, color=TEXT_MUTED)
    footer(s, 14)

    out = Path(__file__).resolve().parent / "export" / "DFU_Portfolio.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
